"""
generate_pptx.py - Generate the complete 28-slide Mobilé Inc. presentation.
Uses pptx_theme.py design system. Reads content from knowledge-base.
Embeds charts from presentation/assets/charts/ and diagrams from assets/diagrams/.
"""

import os
import sys

# Ensure the presentation directory is on the path
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, SCRIPT_DIR)

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from pptx_theme import (
    create_presentation,
    create_title_slide,
    create_section_divider,
    create_content_slide,
    create_table_slide,
    create_diagram_slide,
    create_two_column_slide,
    add_text_box,
    add_slide_number,
    add_image,
    create_styled_table,
    set_cell_text,
    _add_blank_slide,
    _add_title_bar,
    _add_title_with_subtitle,
    add_callout_box,
    add_wave_indicators,
    _add_reference_footnote,
    _add_thin_line,
    _set_slide_bg,
    _rgb_copy,
    # Constants
    WHITE, OFF_WHITE, LIGHT_GRAY_BG,
    TEXT_PRIMARY, TEXT_SECONDARY, TEXT_MUTED,
    WAVE_4G, WAVE_5G, WAVE_AI,
    ACCENT_BLUE, DANGER_RED, AMBER,
    TABLE_HEADER_BG, TABLE_HEADER_TEXT, TABLE_ALT_ROW, TABLE_BORDER,
    FONT_FAMILY, TITLE_SIZE, SUBTITLE_SIZE, BODY_SIZE, SMALL_SIZE,
    TABLE_HEADER_SIZE, TABLE_BODY_SIZE, SLIDE_NUMBER_SIZE,
    SLIDE_WIDTH, SLIDE_HEIGHT,
    MARGIN_LEFT, MARGIN_RIGHT, MARGIN_TOP_CONTENT, MARGIN_BOTTOM,
    TITLE_TOP, TITLE_HEIGHT, CONTENT_WIDTH, CONTENT_HEIGHT,
)

# ---------------------------------------------------------------------------
# Asset paths
# ---------------------------------------------------------------------------
ASSETS = os.path.join(SCRIPT_DIR, "assets")
CHARTS = os.path.join(ASSETS, "charts")
DIAGRAMS = os.path.join(ASSETS, "diagrams")
PHOTOS = os.path.join(ASSETS, "photos")
OUTPUT_DIR = os.path.join(SCRIPT_DIR, "output")


def asset_path(subdir, filename):
    """Return full path to an asset, or None if it doesn't exist."""
    p = os.path.join(ASSETS, subdir, filename)
    return p if os.path.isfile(p) else None


def chart(name):
    return asset_path("charts", name)


def diagram(name):
    return asset_path("diagrams", name)


def photo(name):
    return asset_path("photos", name)


# ---------------------------------------------------------------------------
# Speaker notes helper
# ---------------------------------------------------------------------------
def add_speaker_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------------------
# Build all 28 slides
# ---------------------------------------------------------------------------

def build_presentation():
    prs = create_presentation()
    slide_num = 0

    # ================================================================
    # SLIDE 1: Title Page
    # ================================================================
    slide_num += 1
    slide = create_title_slide(
        prs,
        title="Mobile Inc. - Strategic Analysis 2025",
        subtitle="International Business Strategy with Simulation (MN7002NI)",
        team_members=[
            "Aayush Man Singh (18029779)",
            "Gaurav Dangol (18029811)",
            "Ruchan Jung Sijapati (20048645)",
            "Shaswat Nibha Maharjan",
            "Bishwesh Ram Shrestha (25030247)",
            "Bishan Subedi (25030193)",
        ],
        module_code="MN7002NI",
        date="10 April 2026",
    )
    add_speaker_notes(slide, "Title slide. Module MN7002NI, International Business Strategy with Simulation. Group presentation analyzing Mobile Inc.'s strategic position in 2025.")

    # ================================================================
    # SLIDE 2: Agenda
    # ================================================================
    slide_num += 1
    slide = create_content_slide(
        prs,
        title="Agenda",
        bullets=[
            "Executive Summary",
            "The Three Waves of Value - 4G, 5G, AI",
            "Task 1: Strategy Process - Vision, Mission, Goals, SMART Objectives",
            "Task 2: Internal & External Analysis - VRIO, Value Chain, PESTLE, Porter\u2019s Five Forces, Strategic Group, CSF",
            "Task 3: SWOT Synthesis - Evidence-Based Analysis",
            "Task 4: Strategic Recommendation - Accelerated 5G Entry with Parallel AI R&D",
            "References",
        ],
        slide_number=slide_num,
    )
    add_speaker_notes(slide, "Walk through the structure of our presentation. Four tasks building on each other, with a clear strategic recommendation at the end.")

    # ================================================================
    # SLIDE 3: Executive Summary
    # ================================================================
    slide_num += 1
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Executive Summary")

    # Key metrics callout boxes
    metrics = [
        ("$4.8B", "Annual 4G Revenue", WAVE_4G),
        ("0", "5G Handsets Launched", DANGER_RED),
        ("$90M", "Cash Reserve Floor", AMBER),
        ("2027-28", "AI Device Window", WAVE_AI),
    ]
    box_width = Inches(2.6)
    box_gap = Inches(0.25)
    total_boxes_width = box_width * 4 + box_gap * 3
    start_left = (SLIDE_WIDTH - total_boxes_width) // 2

    for i, (value, label, color) in enumerate(metrics):
        box_left = start_left + i * (box_width + box_gap)
        box_top = MARGIN_TOP_CONTENT

        # Value
        add_text_box(slide, value, box_left, box_top, box_width, Inches(0.5),
                     font_size=Pt(26), color=color, bold=True, alignment=PP_ALIGN.CENTER)
        # Label
        add_text_box(slide, label, box_left, box_top + Inches(0.5), box_width, Inches(0.35),
                     font_size=SMALL_SIZE, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

    # Summary text below metrics
    summary_bullets = [
        "Mobile Inc. is a US smartphone maker competing across USA, Europe, and Asia with 25% market share in each",
        "Post-DOJ intervention (2025) ended tacit price coordination - open competition begins",
        "Three simultaneous technology waves compete for finite resources: defend 4G cash engine, launch 5G urgently, invest in AI R&D for 2027-28",
        "Our recommendation: Accelerated 5G entry via licensing + parallel AI R&D, funded by managed 4G transition",
        "All investments constrained by the $90M cash reserve floor - sequencing is critical",
    ]
    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(2.4), CONTENT_WIDTH, Inches(4.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(summary_bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"\u2022  {bullet}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_PRIMARY
        p.font.name = FONT_FAMILY
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)

    add_slide_number(slide, slide_num)
    add_speaker_notes(slide, "Executive summary: Mobile Inc. faces a three-wave challenge. 4G cash engine funds everything, but 5G product gap is critical. Our recommendation is a sequenced dual-investment strategy.")

    # ================================================================
    # SLIDE 4: Three Waves Overview
    # ================================================================
    slide_num += 1
    three_waves_chart = chart("three_waves_scurve.png")
    if three_waves_chart:
        slide = create_diagram_slide(
            prs,
            title="The Three Waves of Value",
            image_path=three_waves_chart,
            caption="Figure 1: Mobile Inc.\u2019s three technology waves mapped to McKinsey Three Horizons (Baghai, Coley and White, 1999)",
            reference="Source: Case study data; McKinsey Three Horizons framework",
            slide_number=slide_num,
        )
    else:
        slide = create_content_slide(
            prs,
            title="The Three Waves of Value",
            bullets=[
                "Wave 1 (4G LTE): Mature cash engine, ~$4.8B revenue, declining margins post-DOJ",
                "Wave 2 (5G): Urgent - 70%+ urban coverage, Mobile has NO 5G handset yet",
                "Wave 3 (AI-Devices): 2027-28 horizon, R&D investment decisions must happen NOW",
                "All three waves compete for the same finite resource pool",
                "Cannot ignore any wave - each has irreversible consequences for delay",
            ],
            reference="Source: McKinsey Three Horizons (Baghai, Coley and White, 1999)",
            slide_number=slide_num,
        )
    add_speaker_notes(slide, "The three waves are the defining challenge. Wave 1 funds everything. Wave 2 is urgent. Wave 3 window is open now. The central tension: all compete for finite resources, constrained by the $90M cash floor.")

    # ================================================================
    # SLIDE 5: Task 1 - Vision & Mission
    # ================================================================
    slide_num += 1
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Vision & Mission")

    # Vision
    add_text_box(slide, "Vision", MARGIN_LEFT, Inches(1.15), Inches(2), Inches(0.35),
                 font_size=Pt(14), color=WAVE_4G, bold=True)
    add_text_box(slide,
        "\u201cTo lead the global smartphone industry through the 4G-to-5G transition and into the AI-device era by 2028, delivering differentiated technology across the US, European, and Asian markets while sustaining financial resilience through disciplined multi-horizon resource allocation.\u201d",
        MARGIN_LEFT, Inches(1.5), CONTENT_WIDTH, Inches(1.3),
        font_size=Pt(14), color=TEXT_PRIMARY, bold=False)

    # Divider line
    _add_thin_line(slide, MARGIN_LEFT, Inches(2.9), CONTENT_WIDTH, TABLE_BORDER)

    # Mission
    add_text_box(slide, "Mission", MARGIN_LEFT, Inches(3.1), Inches(2), Inches(0.35),
                 font_size=Pt(14), color=WAVE_5G, bold=True)
    add_text_box(slide,
        "\u201cMobile Inc. designs, manufactures, and delivers smartphones that serve the distinct needs of consumers across the US, Europe, and Asia - defending our 4G leadership, launching competitive 5G devices, and investing in AI-integrated device R&D - all while maintaining the financial discipline required to fund three simultaneous technology waves from a single resource base.\u201d",
        MARGIN_LEFT, Inches(3.45), CONTENT_WIDTH, Inches(1.6),
        font_size=Pt(14), color=TEXT_PRIMARY, bold=False)

    _add_reference_footnote(slide, "Aligned with McKinsey Three Horizons (Baghai, Coley and White, 1999)")
    add_slide_number(slide, slide_num)
    add_speaker_notes(slide, "Our vision reflects the three-wave reality: lead through the transition, not just in one technology. The mission is specific to Mobile: what we do, for whom, across which markets, within which constraint.")

    # ================================================================
    # SLIDE 7: Strategic Goals + SMART Objectives
    # ================================================================
    slide_num += 1
    slide = create_table_slide(
        prs,
        title="Strategic Goals & SMART Objectives",
        headers=["Wave", "Strategic Goal", "SMART Objective (Summary)", "Timeline"],
        rows=[
            ["4G (H1)", "Defend 4G revenue base", "Maintain 22% Asia share, 8% cost reduction via Vietnam, revenue floor $4.0B", "Q4 2026"],
            ["5G (H2)", "Achieve competitive 5G entry", "Launch 5G in US + EU by Q2 2026, target 10% US / 8% EU share within 12mo", "Q2 2026-Q2 2027"],
            ["AI (H3)", "Secure AI-device position", "Allocate $150M AI R&D, 2 partnerships, working prototype", "FY2026-Q4 2027"],
            ["Cross", "Maintain financial resilience", "Cash reserves >$120M (30% buffer), max $200M borrowing at 4.8%", "2025-2027"],
            ["Cross", "Optimize global production", "Commission 4 new Vietnam lines, evaluate contract mfg. at 18% cap", "Q1 2026"],
        ],
        reference="SMART criteria (Doran, 1981); Case study data",
        slide_number=slide_num,
        col_widths=[Inches(0.9), Inches(2.5), Inches(6.0), Inches(2.1)],
    )
    add_speaker_notes(slide, "Five strategic goals mapped to the three waves plus cross-cutting financial discipline and production. Each has a SMART objective with measurable targets and deadlines.")

    # ================================================================
    # SLIDE 8: Strategy Hierarchy
    # ================================================================
    slide_num += 1
    goal_hierarchy = diagram("goal_hierarchy_diagram.png")
    if goal_hierarchy:
        slide = create_diagram_slide(
            prs,
            title="Strategy Hierarchy",
            image_path=goal_hierarchy,
            caption="Figure 2: Vision \u2192 Mission \u2192 Goals \u2192 SMART Objectives hierarchy for Mobile Inc.",
            reference="Source: Johnson et al. (2017)",
            slide_number=slide_num,
        )
    else:
        slide = create_content_slide(
            prs,
            title="Strategy Hierarchy",
            bullets=[
                "Vision: Lead global smartphone industry through 4G\u21925G\u2192AI transition by 2028",
                "Mission: Design, manufacture, deliver across three markets with financial discipline",
                "Goals: 5 strategic goals spanning all three waves + cross-cutting concerns",
                "SMART Objectives: Each goal has specific, measurable, time-bound targets",
                "The hierarchy ensures every operational decision traces to the three-wave vision",
            ],
            reference="Source: Johnson et al. (2017) Exploring Strategy",
            slide_number=slide_num,
        )
    add_speaker_notes(slide, "The strategy hierarchy flows from vision down to SMART objectives. Every level references the three-wave structure, ensuring coherence.")

    # ================================================================
    # SLIDE 9: Task 2 Section Divider
    # ================================================================
    slide_num += 1
    slide = create_section_divider(
        prs,
        task_number="Task 2",
        title="Internal & External Analysis",
        description="VRIO, Value Chain, PESTLE, Porter\u2019s Five Forces, Strategic Group, CSF",
        accent_color=WAVE_5G,
    )
    add_slide_number(slide, slide_num)
    add_speaker_notes(slide, "Task 2 covers six frameworks: VRIO, Value Chain, PESTLE, Porter's Five Forces, Strategic Group, and CSF. These provide both internal and external analysis of Mobile Inc.'s strategic position.")

    # ================================================================
    # SLIDE 10: VRIO Analysis (1) - Framework + Table
    # ================================================================
    slide_num += 1
    slide = create_table_slide(
        prs,
        title="VRIO Analysis - Core Resources",
        headers=["Resource / Capability", "V", "R", "I", "O", "Competitive Implication"],
        rows=[
            ["Atlanta plant (12 lines, learning curve)", "Y", "Partial", "Y", "Y", "Temporary advantage"],
            ["Vietnam plant ($6/hr labour)", "Y", "N", "N", "Partial", "Competitive parity"],
            ["4G cash engine ($4.8B)", "Y", "N", "N", "Y", "Competitive parity"],
            ["Brand (25% share, 3 markets)", "Y", "N", "N", "Y", "Competitive parity"],
            ["R&D infrastructure (existing)", "Y", "Partial", "Y", " - ", "Potential advantage"],
            ["Management experience (multi-market)", "Y", "N", "Partial", "Y", "Temporary advantage"],
            ["JIT production capability", "Y", "N", "N", "Y", "Competitive parity"],
        ],
        reference="VRIO framework (Barney, 1991); Grant (2019)",
        slide_number=slide_num,
        col_widths=[Inches(3.2), Inches(0.6), Inches(0.9), Inches(0.6), Inches(0.9), Inches(2.7)],
    )
    add_speaker_notes(slide, "VRIO reveals that Mobile holds NO sustained competitive advantages. Atlanta learning curve is temporary, R&D is potential but unrealized. Most resources are at competitive parity. The critical finding: R&D must be redirected to 5G and AI urgently.")

    # ================================================================
    # SLIDE 11: VRIO Analysis (2) - Key Findings + Diagram
    # ================================================================
    slide_num += 1
    vrio_diag = diagram("vrio_framework_diagram.png")
    slide = create_two_column_slide(
        prs,
        title="VRIO Analysis - Key Findings",
        left_bullets=[
            "No sustained competitive advantages - only temporary (Atlanta) and potential (R&D)",
            "Atlanta's learning-curve efficiencies are real but replicable over time by rivals",
            "Vietnam's cost advantage is location-based, not proprietary - Samsung, Intel, LG also operate there",
            "R&D infrastructure is the most strategically significant resource: IF redirected to 5G and AI, it becomes a potential source of sustained advantage",
            "Critical vulnerability: zero 5G capability despite 70%+ US urban 5G coverage",
        ],
        right_image_path=vrio_diag,
        reference="VRIO framework (Barney, 1991); Case study data",
        slide_number=slide_num,
    )
    add_speaker_notes(slide, "The VRIO verdict is clear: Mobile must convert potential advantages into actual capabilities. The R&D base exists but is not generating 5G or AI outputs.")

    # ================================================================
    # SLIDE 12: Value Chain Analysis (1) - Diagram
    # ================================================================
    slide_num += 1
    vc_diagram = diagram("value_chain_diagram.png")
    if vc_diagram:
        slide = create_diagram_slide(
            prs,
            title="Value Chain Analysis - Primary & Support Activities",
            image_path=vc_diagram,
            caption="Figure 3: Porter\u2019s Value Chain applied to Mobile Inc.",
            reference="Porter (1985); Case study data",
            slide_number=slide_num,
        )
    else:
        slide = create_content_slide(
            prs,
            title="Value Chain Analysis - Overview",
            bullets=[
                "Inbound Logistics: JIT system - no stockpiling, but risky for uncertain 5G demand",
                "Operations: Dual-plant (Atlanta $28/hr premium + Vietnam $6/hr cost-competitive)",
                "Outbound Logistics: Atlanta\u2192US, Vietnam\u2192Asia (zero tariff), EVFTA\u2192Europe",
                "Marketing: Brand capability exists but no 5G product to sell",
                "Technology (R&D): THE most critical support activity - not yet applied to 5G/AI",
            ],
            reference="Porter (1985); Case study data",
            slide_number=slide_num,
        )
    add_speaker_notes(slide, "The value chain shows dual-plant structure as a genuine advantage for global coverage. But the critical bottleneck is R&D: it hasn't been directed toward 5G or AI.")

    # ================================================================
    # SLIDE 13: Value Chain Analysis (2) - Production Comparison
    # ================================================================
    slide_num += 1
    prod_chart = chart("production_cost_comparison.png")
    if prod_chart:
        slide = create_two_column_slide(
            prs,
            title="Value Chain - Manufacturing & Cost Structure",
            left_bullets=[
                "Atlanta: 12 lines, $2.9M/line, $28/hr labour",
                "Vietnam: growing, $3.3M/line, $6/hr labour (78% savings)",
                "New capacity: 2-year lead time from investment to operational",
                "JIT works for 4G but risky for uncertain 5G volumes",
                "Contract manufacturing: 15-20% premium, bridge-only solution",
            ],
            right_image_path=prod_chart,
            reference="Porter (1985); Ferdows (1997); Case data",
            slide_number=slide_num,
        )
    else:
        slide = create_content_slide(
            prs,
            title="Value Chain - Manufacturing & Cost Structure",
            bullets=[
                "Atlanta: 12 lines, $2.9M/line, $28/hr labour - premium production for US/EU",
                "Vietnam: expanding, $3.3M/line, $6/hr labour - cost-competitive for Asia/EU (EVFTA)",
                "Critical bottleneck: 2-year lead time for new production capacity",
                "Capacity decisions in 2025 determine 2027 capability (5G peak + AI launch)",
                "Contract manufacturing (15-20% premium) is the only bridge - not a long-term strategy",
            ],
            reference="Porter (1985); Ferdows (1997); Case study data",
            slide_number=slide_num,
        )
    add_speaker_notes(slide, "Manufacturing comparison highlights the dual-plant advantage. The 2-year lead time is the critical constraint: capacity decisions cannot be deferred.")

    # ================================================================
    # SLIDE 14: PESTLE Analysis (1) - Diagram
    # ================================================================
    slide_num += 1
    pestle_diag = diagram("pestle_diagram.png")
    if pestle_diag:
        slide = create_diagram_slide(
            prs,
            title="PESTLE Analysis - Macro-Environmental Forces",
            image_path=pestle_diag,
            caption="Figure 4: PESTLE factors affecting Mobile Inc. across three markets",
            reference="Johnson et al. (2017); Case study data",
            slide_number=slide_num,
        )
    else:
        slide = create_content_slide(
            prs,
            title="PESTLE Analysis - Overview",
            bullets=[
                "Political: DOJ antitrust intervention (most impactful), US-Vietnam trade tensions",
                "Economic: $90M cash floor, 4.8%/7.8% borrowing, Asian growth 22-30%",
                "Social: 38% US 5G adoption, AI use cases exceeding forecasts, EU sustainability",
                "Technological: 70%+ 5G coverage ready, AI on-device chips emerging, 2-year production lag",
                "Legal: DOJ enforcement, EVFTA enabling Vietnam\u2192EU, IP protection for in-house R&D",
                "Environmental: EU e-waste regulations, ESG expectations, sustainability differentiation",
            ],
            reference="Johnson et al. (2017); GSMA (2024); Case study data",
            slide_number=slide_num,
        )
    add_speaker_notes(slide, "PESTLE reveals the DOJ intervention as the single most impactful force. Technology readiness for 5G is not the barrier - Mobile's own product gap is.")

    # ================================================================
    # SLIDE 15: PESTLE Analysis (2) - Key Findings Table
    # ================================================================
    slide_num += 1
    slide = create_table_slide(
        prs,
        title="PESTLE - Key Findings by Factor",
        headers=["Factor", "Key Finding for Mobile Inc.", "Wave Impact"],
        rows=[
            ["Political", "DOJ ended tacit price coordination; open competition begins", "All waves"],
            ["Economic", "$90M cash floor constrains simultaneous investment", "All waves"],
            ["Social", "38% US / 28% EU already own 5G; demand exceeding forecasts", "Waves 2 & 3"],
            ["Technological", "5G infrastructure ready; AI chips emerging (Apple, Samsung active)", "Waves 2 & 3"],
            ["Legal", "EVFTA enables Vietnam\u2192EU; IP protection favours in-house R&D", "Waves 1 & 2"],
            ["Environmental", "EU sustainability = differentiation opportunity", "Waves 1 & 2"],
        ],
        reference="GSMA (2024); Ericsson (2024); Case study data",
        slide_number=slide_num,
    )
    add_speaker_notes(slide, "The PESTLE table summarises each factor with its specific impact on Mobile and which wave is affected. The DOJ intervention and 5G readiness are the dominant forces.")

    # ================================================================
    # SLIDE 16: Porter's Five Forces (1) - Diagram
    # ================================================================
    slide_num += 1
    five_forces = diagram("five_forces_diagram.png")
    if five_forces:
        slide = create_diagram_slide(
            prs,
            title="Porter\u2019s Five Forces - Industry Structure",
            image_path=five_forces,
            caption="Figure 5: Five Forces analysis of the smartphone industry post-DOJ",
            reference="Porter (1980); Case study data",
            slide_number=slide_num,
        )
    else:
        slide = create_content_slide(
            prs,
            title="Porter\u2019s Five Forces - Overview",
            bullets=[
                "Competitive Rivalry: HIGH - 4 equal firms (25% each), post-DOJ open competition",
                "Threat of New Entrants: LOW - high capital requirements, economies of scale",
                "Buyer Power: MODERATE to HIGH - low switching costs, rising price transparency",
                "Supplier Power: MODERATE - 5G licensing creates dependency; AI chips emerging bottleneck",
                "Threat of Substitutes: LOW - no real substitute for smartphones in 2025",
            ],
            reference="Porter (1980); Case study data",
            slide_number=slide_num,
        )
    add_speaker_notes(slide, "Porter's Five Forces: rivalry is the dominant force. The symmetric oligopoly (4 firms at 25% each) makes first-mover advantage in 5G the primary mechanism for breaking from the pack.")

    # ================================================================
    # SLIDE 17: Porter's Five Forces (2) - Key Findings
    # ================================================================
    slide_num += 1
    slide = create_table_slide(
        prs,
        title="Porter\u2019s Five Forces - Assessment",
        headers=["Force", "Intensity", "Key Driver", "Impact on Mobile"],
        rows=[
            ["Competitive Rivalry", "HIGH", "4 equal firms, post-DOJ", "Race for 5G differentiation"],
            ["New Entrants", "LOW", "High capital + scale barriers", "Focus on existing rivals"],
            ["Buyer Power", "MOD-HIGH", "Low switching costs", "5G brand loyalty mitigates"],
            ["Supplier Power", "MODERATE", "5G licensing, AI chips", "In-house R&D reduces dependency"],
            ["Substitutes", "LOW", "No smartphone replacement", "Risk is within-category (4G\u21925G)"],
        ],
        reference="Porter (1980); Case study data",
        slide_number=slide_num,
    )
    add_speaker_notes(slide, "The symmetric four-firm structure means any competitor's 5G launch immediately erodes Mobile's position. First-mover advantage is especially durable in symmetric oligopolies.")

    # ================================================================
    # SLIDE 18: Strategic Group Analysis
    # ================================================================
    slide_num += 1
    sg_chart = chart("strategic_group_map.png")
    if sg_chart:
        slide = create_diagram_slide(
            prs,
            title="Strategic Group Analysis - Migration Paths",
            image_path=sg_chart,
            caption="Figure 6: Strategic group positions and Mobile\u2019s target trajectory (2025-2028)",
            reference="Johnson et al. (2017); Case study data",
            slide_number=slide_num,
        )
    else:
        slide = create_content_slide(
            prs,
            title="Strategic Group Analysis",
            bullets=[
                "Current: All 4 firms clustered as 'Global 4G Incumbents' - 25% share each",
                "Group A \u2018 5G First-Movers\u2019: capture premium pricing + loyalty (target for Mobile)",
                "Group B \u20184G Defenders\u2019: strong in Asian mass-market but technology followers",
                "Group C \u2018Technology Leaders\u2019: heavy 5G + AI R&D investment, future dominance",
                "Mobile\u2019s trajectory: Global 4G Incumbent \u2192 5G First-Mover \u2192 Technology Leader",
            ],
            reference="Johnson et al. (2017); Case study data",
            slide_number=slide_num,
        )
    add_speaker_notes(slide, "Strategic groups are currently identical but will diverge. Mobile should target Group A (5G First-Mover) immediately while investing to join Group C (Technology Leader). The window for group migration is time-limited.")

    # ================================================================
    # SLIDE 19: CSF Analysis
    # ================================================================
    slide_num += 1
    csf_chart = chart("csf_gap_analysis.png")
    if csf_chart:
        slide = create_two_column_slide(
            prs,
            title="Critical Success Factors - Gap Analysis",
            left_bullets=[
                "Speed to 5G market: CRITICAL gap (most urgent)",
                "AI R&D commitment: CRITICAL gap (window closing)",
                "4G cash defense: GROWING gap (margins compressing)",
                "Financial discipline: EMERGING gap ($90M floor pressure)",
                "Production capacity: MODERATE gap (2-year lag)",
                "JIT capability: NO gap (established)",
            ],
            right_image_path=csf_chart,
            reference="CSF analysis; Case study data",
            slide_number=slide_num,
        )
    else:
        slide = create_content_slide(
            prs,
            title="Critical Success Factors - Priority Ranking",
            bullets=[
                "1. Speed to 5G market (Wave 2) - most time-sensitive, largest gap",
                "2. AI R&D commitment (Wave 3) - window open but closing; irreversible",
                "3. 4G cash defense (Wave 1) - funds everything else",
                "4. Financial discipline (Cross-wave) - $90M floor shapes all decisions",
                "5. Production capacity (Cross-wave) - 2-year lead time makes this urgent",
            ],
            reference="CSF analysis; Case study data",
            slide_number=slide_num,
        )
    add_speaker_notes(slide, "CSF analysis confirms 5G speed-to-market as the single most critical success factor. The gap analysis shows two critical, two growing, one moderate, and one resolved gap.")

    # ================================================================
    # SLIDE 19: SWOT Matrix (Task 3)
    # ================================================================
    slide_num += 1
    swot_diag = diagram("swot_matrix_diagram.png")
    if swot_diag:
        slide = create_diagram_slide(
            prs,
            title="SWOT Matrix - Mobile Inc.",
            image_path=swot_diag,
            caption="All SWOT points traced to Task 2 framework evidence",
            reference="SWOT synthesis drawing from all Task 2 frameworks",
            slide_number=slide_num,
        )
    else:
        slide = create_table_slide(
            prs,
            title="SWOT Matrix - Mobile Inc.",
            headers=["", "Positive", "Negative"],
            rows=[
                ["Internal",
                 "S1: Dual-plant mfg.\nS2: $4.8B cash engine\nS3: R&D + global mgmt\nS4: JIT efficiency",
                 "W1: No 5G handset\nW2: 100% 4G dependency\nW3: Vietnam immaturity\nW4: R&D not directed"],
                ["External",
                 "O1: 5G ready\nO2: AI window open\nO3: Asia growth\nO4: EVFTA",
                 "T1: Post-DOJ price war\nT2: Rival 5G first-mover\nT3: Tariff risk\nT4: AI competitors"],
            ],
            reference="All SWOT points traced to Task 2 framework evidence",
            slide_number=slide_num,
        )
    add_speaker_notes(slide, "Every SWOT point traces to a specific Task 2 framework finding. No new information is introduced here. The dominant pattern: sufficient capabilities, misaligned with external reality.")

    # ================================================================
    # SLIDE 20: SWOT Connections & Implications
    # ================================================================
    slide_num += 1
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "SWOT - Connections & Strategic Implications")

    swot_bullets = [
        "W1 \u2192 O1: No 5G handset + 5G infrastructure ready = URGENT launch opportunity",
        "W4 \u2192 O2: R&D not directed + AI window open = invest NOW or forfeit permanently",
        "S2 \u2192 T1: 4G cash engine + post-DOJ price war = defend revenue to fund Waves 2 & 3",
        "S1 \u2192 O3: Dual-plant structure + Asian growth = Vietnam expansion for cost-competitive production",
        "Financial discipline ($90M floor) means investments must be sequenced, not simultaneous at maximum",
    ]
    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP_CONTENT, CONTENT_WIDTH, Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(swot_bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"\u2022  {bullet}"
        p.font.size = BODY_SIZE
        p.font.color.rgb = _rgb_copy(TEXT_PRIMARY)
        p.font.name = FONT_FAMILY
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)

    # Callout box at bottom
    add_callout_box(
        slide,
        "Key Insight: Capabilities are sufficient but MISALIGNED with external reality. "
        "The strategic task is realignment, not reinvention.",
        left=MARGIN_LEFT,
        top=Inches(5.5),
        width=CONTENT_WIDTH,
        height=Inches(0.65),
        accent_color=WAVE_5G,
    )

    _add_reference_footnote(slide, "SWOT synthesis drawing from all Task 2 frameworks; feeds into Task 4 recommendation")
    add_slide_number(slide, slide_num)
    add_speaker_notes(slide, "SWOT connections and implications combined. The dominant move: close 5G gap (W1-O1), commit to AI R&D (W4-O2), fund from 4G defense (S2-T1). Realignment, not transformation.")

    # ================================================================
    # SLIDE 24: Task 4 Section Divider
    # ================================================================
    slide_num += 1
    slide = create_section_divider(
        prs,
        task_number="Task 4",
        title="Strategic Recommendation",
        description="Accelerated 5G Entry with Parallel AI R&D, Funded by Managed 4G Transition",
        accent_color=WAVE_AI,
    )
    add_slide_number(slide, slide_num)
    add_speaker_notes(slide, "Task 4: our strategic recommendation. A sequenced dual-investment strategy that directly addresses every finding from Tasks 1-3.")

    # ================================================================
    # SLIDE 25: Strategic Choice
    # ================================================================
    slide_num += 1
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "The Strategic Choice")

    choice_bullets = [
        "Accelerated 5G entry via technology licensing (6-9 months to market)",
        "Parallel in-house AI R&D to secure long-term technology leadership",
        "Managed 4G transition: defend cash generation, accept controlled decline",
        "All funded from existing 4G revenue + disciplined borrowing (4.8% rate)",
        "Key principle: sequence investments, do not attempt everything at maximum simultaneously",
    ]
    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP_CONTENT, CONTENT_WIDTH, Inches(3.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(choice_bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"\u2022  {bullet}"
        p.font.size = BODY_SIZE
        p.font.color.rgb = _rgb_copy(TEXT_PRIMARY)
        p.font.name = FONT_FAMILY
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)

    # Three-column comparison boxes at bottom
    col_w = Inches(3.7)
    col_gap = Inches(0.2)
    total_w = 3 * col_w + 2 * col_gap
    start_x = (SLIDE_WIDTH - total_w) // 2
    box_top = Inches(5.0)
    box_h = Inches(0.85)
    wave_data = [
        ("4G Defense", "Defend $4.0B floor, 8% cost cut via Vietnam", WAVE_4G),
        ("5G Launch", "License + launch Q2 2026, 10% US share target", WAVE_5G),
        ("AI R&D", "$150M investment, prototype by Q4 2027", WAVE_AI),
    ]
    for i, (label, desc, color) in enumerate(wave_data):
        x = start_x + i * (col_w + col_gap)
        add_callout_box(slide, f"{label}: {desc}", x, box_top, col_w, box_h, color)

    _add_reference_footnote(slide, "Aligned with Tasks 1-3; Teece, Pisano and Shuen (1997); Doz and Kosonen (2010)")
    add_slide_number(slide, slide_num)
    add_speaker_notes(slide, "Our recommendation is a sequenced dual-investment strategy. License 5G for speed, build AI capability for the long term, fund everything from managed 4G cash flows. This directly addresses every SWOT finding.")

    # ================================================================
    # SLIDE 26: 4G Action Plan
    # ================================================================
    slide_num += 1
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)
    content_start = _add_title_with_subtitle(
        slide,
        title="Wave 1: 4G Action Plan",
        subtitle="Defend, Extract, Don't Over-Invest",
        subtitle_color=WAVE_4G,
    )

    bullets_4g = [
        "Maintain 4G production: Atlanta (premium US/EU) + expand Vietnam (cost-competitive Asia)",
        "Target 8% per-unit cost reduction via Vietnam scale-up",
        "Accept managed share erosion: 25% -> 22% floor - deliberate trade-off, not failure",
        "Price aggressively in Asian mass-market using Vietnam cost advantage",
        "Optimize, don't maximize: extract cash to fund Waves 2 and 3",
    ]
    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(1.6), CONTENT_WIDTH, Inches(3.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets_4g):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"\u2022  {bullet}"
        p.font.size = BODY_SIZE
        p.font.color.rgb = _rgb_copy(TEXT_PRIMARY)
        p.font.name = FONT_FAMILY
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)

    add_callout_box(
        slide,
        "Key metric: Maintain revenue floor $4.0B while reducing costs 8% via Vietnam",
        left=MARGIN_LEFT,
        top=Inches(5.5),
        width=CONTENT_WIDTH,
        height=Inches(0.6),
        accent_color=WAVE_4G,
    )

    _add_reference_footnote(slide, "VRIO (Barney, 1991); Value Chain (Porter, 1985); Case study data")
    add_slide_number(slide, slide_num)
    add_speaker_notes(slide, "4G strategy is about defense and extraction. Over-investing in 4G at the expense of 5G launch would be the single biggest strategic mistake. Accept decline, extract cash.")

    # ================================================================
    # SLIDE 27: 5G Action Plan
    # ================================================================
    slide_num += 1
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)
    content_start = _add_title_with_subtitle(
        slide,
        title="Wave 2: 5G Action Plan",
        subtitle="License to Launch Fast, Build IP in Parallel",
        subtitle_color=WAVE_5G,
    )

    bullets_5g = [
        "License 5G chipset technology for immediate market entry (6-9 months)",
        "Target US and European premium segments first (brand-conscious, 5G-ready infrastructure)",
        "Launch by Q2 2026 - every quarter of delay cedes first-mover advantage",
        "Begin in-house 5G R&D immediately - second-gen product (12-18 months) builds proprietary IP",
        "Use contract manufacturing (15-20% premium) to bridge production gap",
        "Target: 10% US / 8% EU 5G market share within 12 months of launch",
    ]
    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(1.6), CONTENT_WIDTH, Inches(3.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets_5g):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"\u2022  {bullet}"
        p.font.size = BODY_SIZE
        p.font.color.rgb = _rgb_copy(TEXT_PRIMARY)
        p.font.name = FONT_FAMILY
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)

    add_callout_box(
        slide,
        "Key metric: 10% US / 8% EU 5G market share within 12 months of Q2 2026 launch",
        left=MARGIN_LEFT,
        top=Inches(5.5),
        width=CONTENT_WIDTH,
        height=Inches(0.6),
        accent_color=WAVE_5G,
    )

    _add_reference_footnote(slide, "Strategic Group analysis; CSF priority ranking; Case study data")
    add_slide_number(slide, slide_num)
    add_speaker_notes(slide, "The licensing-first approach is critical: speed beats perfection. In-house R&D runs in parallel for second-gen products that build proprietary IP and reduce supplier dependency.")

    # ================================================================
    # SLIDE 28: AI Action Plan
    # ================================================================
    slide_num += 1
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)
    content_start = _add_title_with_subtitle(
        slide,
        title="Wave 3: AI Action Plan",
        subtitle="Invest Now, Small but Committed",
        subtitle_color=WAVE_AI,
    )

    bullets_ai = [
        "Allocate $150M to AI-integrated device R&D over FY2025-2026",
        "Establish dedicated AI device team (separate from 5G team)",
        "Secure 2 strategic technology partnerships (chip design + AI software)",
        "Target working prototype by Q4 2027",
        "Focus: on-device AI for photography, health monitoring, personalised assistants",
        "Commercial launch 2028 - alongside or ahead of competitors",
    ]
    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(1.6), CONTENT_WIDTH, Inches(3.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets_ai):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"\u2022  {bullet}"
        p.font.size = BODY_SIZE
        p.font.color.rgb = _rgb_copy(TEXT_PRIMARY)
        p.font.name = FONT_FAMILY
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)

    add_callout_box(
        slide,
        "Key metric: $150M R&D investment, 2 partnerships, working prototype by Q4 2027",
        left=MARGIN_LEFT,
        top=Inches(5.5),
        width=CONTENT_WIDTH,
        height=Inches(0.6),
        accent_color=WAVE_AI,
    )

    _add_reference_footnote(slide, "PESTLE (technology); Gartner (2024); Qualcomm (2024); Case study data")
    add_slide_number(slide, slide_num)
    add_speaker_notes(slide, "AI investment is small relative to 4G revenue but the commitment must happen now. Waiting until 5G is profitable before investing = permanently forfeiting technology leadership.")

    # ================================================================
    # SLIDE 29: Market Prioritisation + Trade-offs
    # ================================================================
    slide_num += 1
    market_chart = chart("market_comparison.png")

    # Build this as a two-column: table on left, chart on right
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "Market Prioritisation & Investment Allocation")

    # Market table
    table_data = [
        ["Market", "4G (H1)", "5G (H2)", "AI (H3)"],
        ["USA", "Maintain premium", "PRIMARY launch", "Premium AI target"],
        ["Europe", "Atlanta + Vietnam", "SECONDARY launch", "Feature-rich AI"],
        ["Asia", "PRIMARY defense", "Selective premium", "Later-phase rollout"],
    ]
    table_left = MARGIN_LEFT
    table_top = MARGIN_TOP_CONTENT
    table_width = Inches(7.0)
    table_height = Inches(2.0)

    table = create_styled_table(slide, 4, 4, table_left, table_top, table_width, table_height)
    for j, hdr in enumerate(table_data[0]):
        set_cell_text(table.cell(0, j), hdr, font_size=TABLE_HEADER_SIZE,
                      color=TABLE_HEADER_TEXT, bold=True, alignment=PP_ALIGN.LEFT)
    for i in range(1, 4):
        for j in range(4):
            color = TEXT_PRIMARY
            if j == 1:
                color = WAVE_4G
            elif j == 2:
                color = WAVE_5G
            elif j == 3:
                color = WAVE_AI
            set_cell_text(table.cell(i, j), table_data[i][j], font_size=TABLE_BODY_SIZE,
                          color=color if j > 0 else TEXT_PRIMARY, bold=(j > 0))

    # Chart on the right (if available)
    if market_chart:
        add_image(slide, market_chart, Inches(7.5), MARGIN_TOP_CONTENT, width=Inches(5.0))

    # Trade-off bullets below
    tradeoffs = [
        "4G share decline is deliberate - resources freed fund Waves 2 & 3",
        "Licensing 5G = speed at the cost of short-term supplier dependency",
        "AI investment has uncertain returns 3-4 years out - mitigated by partnerships",
        "Cash floor ($90M) means we cannot do everything at maximum scale simultaneously",
    ]
    txBox = slide.shapes.add_textbox(MARGIN_LEFT, Inches(3.6), CONTENT_WIDTH, Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(tradeoffs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"\u2022  {bullet}"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_PRIMARY
        p.font.name = FONT_FAMILY
        p.space_after = Pt(5)

    _add_reference_footnote(slide, "Source: Case study data; GSMA (2024); Hill (2021)")
    add_slide_number(slide, slide_num)
    add_speaker_notes(slide, "Market prioritisation varies by wave. USA leads 5G, Asia leads 4G defense, Europe bridges both. Trade-offs are addressed honestly - this is a constrained optimisation, not a wish list.")

    # ================================================================
    # SLIDE 30: Investment + Risk (with pie chart)
    # ================================================================
    slide_num += 1
    inv_chart = chart("investment_allocation.png")
    if inv_chart:
        slide = create_two_column_slide(
            prs,
            title="Trade-offs, Risk & Investment Allocation",
            left_bullets=[
                "Cash reserves: maintain $120M minimum (30% above $90M floor)",
                "Borrowing: max $200M at 4.8%, zero emergency borrowing (7.8%)",
                "Risk 1: 5G licensing dependency \u2192 mitigated by parallel in-house R&D",
                "Risk 2: 4G decline faster than expected \u2192 mitigated by Vietnam cost reduction",
                "Risk 3: AI R&D uncertain \u2192 mitigated by partnership risk-sharing",
                "Risk 4: Vietnam tariffs \u2192 mitigated by maintaining Atlanta for US market",
            ],
            right_image_path=inv_chart,
            reference="Case study data; Financial constraints from case brief",
            slide_number=slide_num,
        )
    else:
        slide = create_content_slide(
            prs,
            title="Trade-offs, Risk & Investment Allocation",
            bullets=[
                "Cash reserves: maintain $120M minimum (30% buffer above $90M floor)",
                "Borrowing: max $200M at standard 4.8% rate, zero emergency borrowing",
                "Investment split: ~50% 4G defense, ~30% 5G launch, ~15% AI R&D, ~5% cash buffer",
                "Key risks: 5G licensing dependency, 4G decline speed, AI uncertainty, Vietnam tariffs",
                "Each risk has a specific mitigation strategy - no risk is ignored",
            ],
            reference="Case study data; Financial constraints",
            slide_number=slide_num,
        )
    add_speaker_notes(slide, "Financial discipline is non-negotiable. Every risk has a mitigation. The investment allocation reflects the sequencing principle: fund 4G defense first, then 5G launch, then AI R&D.")

    # ================================================================
    # SLIDE 31 (28 in spec): References
    # ================================================================
    slide_num += 1
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, "References")

    references_col1 = [
        "Baghai, M., Coley, S. and White, D. (1999) The Alchemy of Growth. New York: Perseus Books.",
        "Barney, J.B. (1991) 'Firm Resources and Sustained Competitive Advantage', Journal of Management, 17(1), pp. 99-120.",
        "Cavusgil, S.T., Knight, G. and Riesenberger, J.R. (2017) International Business. 4th edn. Harlow: Pearson.",
        "Deloitte (2024) Global Mobile Consumer Trends. Available at: deloitte.com (Accessed: April 2026).",
        "Doz, Y.L. and Kosonen, M. (2010) 'Embedding Strategic Agility', Long Range Planning, 43(2-3), pp. 370-382.",
        "Ericsson (2024) Ericsson Mobility Report. Available at: ericsson.com (Accessed: April 2026).",
        "European Commission (2020) EU-Vietnam FTA Summary. Available at: trade.ec.europa.eu (Accessed: April 2026).",
        "Ferdows, K. (1997) 'Making the Most of Foreign Factories', Harvard Business Review, 75(2), pp. 73-88.",
        "Gartner (2024) Hype Cycle for Artificial Intelligence. Stamford: Gartner.",
        "Ghemawat, P. (2007) Redefining Global Strategy. Boston: Harvard Business School Press.",
    ]

    references_col2 = [
        "Grant, R.M. (2019) Contemporary Strategy Analysis. 10th edn. Chichester: Wiley.",
        "GSMA (2024) The Mobile Economy 2024. Available at: gsma.com (Accessed: April 2026).",
        "Hill, C.W.L. (2021) International Business. 13th edn. New York: McGraw-Hill.",
        "IDC (2024) Worldwide Smartphone Market Share. Available at: idc.com (Accessed: April 2026).",
        "ITU (2024) Measuring Digital Development. Geneva: ITU.",
        "Johnson, G. et al. (2017) Exploring Strategy. 11th edn. Harlow: Pearson.",
        "McKinsey (2023) The State of 5G. Available at: mckinsey.com (Accessed: April 2026).",
        "Peng, M.W. (2019) Global Strategy. 4th edn. Mason: Cengage.",
        "Porter, M.E. (1980) Competitive Strategy. New York: Free Press.",
        "Porter, M.E. (1985) Competitive Advantage. New York: Free Press.",
        "Qualcomm (2024) The Future of AI on Device. Available at: qualcomm.com (Accessed: April 2026).",
        "Shingo, S. (1989) A Study of the Toyota Production System. Cambridge: Productivity Press.",
        "Teece, D.J., Pisano, G. and Shuen, A. (1997) 'Dynamic Capabilities', SMJ, 18(7), pp. 509-533.",
        "WTO (2024) World Trade Statistical Review 2024. Geneva: WTO.",
    ]

    # Two-column references
    gap = Inches(0.3)
    col_width = int((CONTENT_WIDTH - gap) / 2)

    txBox_left = slide.shapes.add_textbox(MARGIN_LEFT, MARGIN_TOP_CONTENT, col_width, Inches(5.5))
    tf_left = txBox_left.text_frame
    tf_left.word_wrap = True
    for idx, ref in enumerate(references_col1):
        p = tf_left.paragraphs[0] if idx == 0 else tf_left.add_paragraph()
        p.text = ref
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_PRIMARY
        p.font.name = FONT_FAMILY
        p.space_after = Pt(4)

    txBox_right = slide.shapes.add_textbox(MARGIN_LEFT + col_width + gap, MARGIN_TOP_CONTENT, col_width, Inches(5.5))
    tf_right = txBox_right.text_frame
    tf_right.word_wrap = True
    for idx, ref in enumerate(references_col2):
        p = tf_right.paragraphs[0] if idx == 0 else tf_right.add_paragraph()
        p.text = ref
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_PRIMARY
        p.font.name = FONT_FAMILY
        p.space_after = Pt(4)

    add_slide_number(slide, slide_num)
    add_speaker_notes(slide, "Harvard-formatted references. Mix of academic texts (Barney, Porter, Teece, Johnson) and industry reports (GSMA, Ericsson, IDC, Gartner) demonstrates extensive reading beyond course content.")

    return prs, slide_num


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print("Generating Mobile Inc. Strategic Analysis PPTX...")

    prs, total_slides = build_presentation()

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    output_path = os.path.join(OUTPUT_DIR, "mobile_inc_strategic_analysis.pptx")
    prs.save(output_path)

    file_size = os.path.getsize(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"  Total slides: {total_slides}")
    print(f"  File size: {file_size:,} bytes ({file_size / 1024:.1f} KB)")
    print(f"  Dimensions: {prs.slide_width} x {prs.slide_height} EMU (16:9)")

    if file_size < 100 * 1024:
        print("  WARNING: File size < 100KB. May need more embedded content.")
    else:
        print("  File size check: PASSED (>100KB)")

    if total_slides < 28:
        print(f"  WARNING: Only {total_slides} slides. Target is 28.")
    elif total_slides > 30:
        print(f"  WARNING: {total_slides} slides exceeds 30-slide limit!")
    else:
        print(f"  Slide count check: PASSED ({total_slides} slides, within 28-30 range)")


if __name__ == "__main__":
    main()
