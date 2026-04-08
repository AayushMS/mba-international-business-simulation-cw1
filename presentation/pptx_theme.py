"""
pptx_theme.py — Reusable python-pptx design system for the Mobilé Inc.
MBA coursework presentation (MN7002NI, 28 slides, 16:9).

Design principles:
  - Light, minimal, student-made aesthetic (Calibri everywhere)
  - NO drop shadows, NO gradients, NO animations
  - Three-wave color coding: 4G blue, 5G green, AI purple
  - Anti-AI: intentionally plain tables, small slide numbers, breathing room

Compatible with python-pptx 1.0.2.
"""

from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Color Constants
# ---------------------------------------------------------------------------

# Backgrounds
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xFA, 0xFB, 0xFC)
LIGHT_GRAY_BG = RGBColor(0xF8, 0xF9, 0xFA)

# Text
TEXT_PRIMARY = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_SECONDARY = RGBColor(0x6B, 0x72, 0x80)
TEXT_MUTED = RGBColor(0x9C, 0xA3, 0xAF)

# Wave Colors
WAVE_4G = RGBColor(0x3B, 0x82, 0xF6)   # Blue
WAVE_5G = RGBColor(0x10, 0xB9, 0x81)   # Green
WAVE_AI = RGBColor(0x8B, 0x5C, 0xF6)   # Purple

# Accents
ACCENT_BLUE = RGBColor(0x0E, 0xA5, 0xE9)
DANGER_RED = RGBColor(0xEF, 0x44, 0x44)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

# Table
TABLE_HEADER_BG = RGBColor(0x1A, 0x1A, 0x2E)
TABLE_HEADER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_ALT_ROW = RGBColor(0xF8, 0xF9, 0xFA)
TABLE_BORDER = RGBColor(0xE5, 0xE7, 0xEB)

# ---------------------------------------------------------------------------
# Font Constants
# ---------------------------------------------------------------------------

FONT_FAMILY = "Calibri"
TITLE_SIZE = Pt(28)
SUBTITLE_SIZE = Pt(18)
BODY_SIZE = Pt(16)
SMALL_SIZE = Pt(12)
TABLE_HEADER_SIZE = Pt(14)
TABLE_BODY_SIZE = Pt(12)
SLIDE_NUMBER_SIZE = Pt(10)

# ---------------------------------------------------------------------------
# Slide Dimensions (16:9 widescreen)
# ---------------------------------------------------------------------------

SLIDE_WIDTH = Emu(12192000)   # 13.333"
SLIDE_HEIGHT = Emu(6858000)   # 7.5"

# Content margins
MARGIN_LEFT = Inches(0.8)
MARGIN_RIGHT = Inches(0.8)
MARGIN_TOP_CONTENT = Inches(1.2)   # below title area
MARGIN_BOTTOM = Inches(0.6)

# Title area
TITLE_TOP = Inches(0.4)
TITLE_HEIGHT = Inches(0.8)

# Derived: usable content width
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
CONTENT_HEIGHT = SLIDE_HEIGHT - MARGIN_TOP_CONTENT - MARGIN_BOTTOM


# ---------------------------------------------------------------------------
# Low-Level Helpers
# ---------------------------------------------------------------------------

def _rgb_copy(color: RGBColor) -> RGBColor:
    """Return a fresh RGBColor copy (avoids shared-reference issues)."""
    return RGBColor(color[0], color[1], color[2])


def _set_slide_bg(slide, color: RGBColor) -> None:
    """Set a solid fill background on *slide* — no gradients."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb_copy(color)


def add_text_box(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    font_size=BODY_SIZE,
    color=TEXT_PRIMARY,
    bold: bool = False,
    alignment=PP_ALIGN.LEFT,
    font_name: str = FONT_FAMILY,
    vertical_anchor=MSO_ANCHOR.TOP,
):
    """Add a text box with configurable font, color, size, and alignment."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    # Vertical alignment inside box
    try:
        tf.paragraphs[0].alignment = alignment
    except Exception:
        pass

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = _rgb_copy(color)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment

    # Vertical anchor on the text frame
    txBox.text_frame.auto_size = None
    try:
        txBox.text_frame._txBody.bodyPr.set("anchor", {
            MSO_ANCHOR.TOP: "t",
            MSO_ANCHOR.MIDDLE: "ctr",
            MSO_ANCHOR.BOTTOM: "b",
        }.get(vertical_anchor, "t"))
    except Exception:
        pass

    return txBox


def add_image(slide, image_path: str, left, top, width=None, height=None):
    """Embed an image, maintaining aspect ratio when only one dimension given."""
    if width is not None and height is not None:
        return slide.shapes.add_picture(image_path, left, top, width, height)
    elif width is not None:
        return slide.shapes.add_picture(image_path, left, top, width=width)
    elif height is not None:
        return slide.shapes.add_picture(image_path, left, top, height=height)
    else:
        return slide.shapes.add_picture(image_path, left, top)


def set_cell_text(
    cell,
    text: str,
    font_size=TABLE_BODY_SIZE,
    color=TEXT_PRIMARY,
    bold: bool = False,
    alignment=PP_ALIGN.LEFT,
):
    """Style an individual table cell's text and remove default margins."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = _rgb_copy(color)
    p.font.bold = bold
    p.font.name = FONT_FAMILY
    p.alignment = alignment
    # Tighten internal margins
    cell.text_frame.word_wrap = True
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)


def _set_cell_border(cell, color: RGBColor, width: Pt = Pt(0.5)):
    """Set thin borders on a table cell via XML manipulation."""
    from pptx.oxml.ns import qn
    import copy

    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    for edge in ("lnL", "lnR", "lnT", "lnB"):
        ln = tcPr.find(qn(f"a:{edge}"))
        if ln is not None:
            tcPr.remove(ln)

        from lxml import etree
        ln = etree.SubElement(tcPr, qn(f"a:{edge}"))
        ln.set("w", str(int(width)))
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        solidFill = etree.SubElement(ln, qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")


def _set_cell_fill(cell, color: RGBColor):
    """Set solid fill on a table cell."""
    cell.fill.solid()
    cell.fill.fore_color.rgb = _rgb_copy(color)


# ---------------------------------------------------------------------------
# Slide Number
# ---------------------------------------------------------------------------

def add_slide_number(slide, number: int) -> None:
    """Add a small gray slide number to the bottom-right corner."""
    add_text_box(
        slide,
        str(number),
        left=SLIDE_WIDTH - Inches(0.8),
        top=SLIDE_HEIGHT - Inches(0.45),
        width=Inches(0.5),
        height=Inches(0.3),
        font_size=SLIDE_NUMBER_SIZE,
        color=TEXT_MUTED,
        bold=False,
        alignment=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# Styled Table Builder
# ---------------------------------------------------------------------------

def create_styled_table(slide, rows: int, cols: int, left, top, width, height):
    """
    Add a table shape and return the Table object with:
      - Dark header row (TABLE_HEADER_BG + white bold text)
      - Alternating white / light-gray data rows
      - Thin borders
    Caller is responsible for populating cell text via set_cell_text.
    """
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Disable built-in banding so we control colours manually
    table.first_row = False
    table.horz_banding = False
    table.vert_banding = False

    for col_idx in range(cols):
        # Header row
        hdr_cell = table.cell(0, col_idx)
        _set_cell_fill(hdr_cell, TABLE_HEADER_BG)
        _set_cell_border(hdr_cell, TABLE_BORDER)

        # Data rows: alternating colours
        for row_idx in range(1, rows):
            data_cell = table.cell(row_idx, col_idx)
            bg = WHITE if row_idx % 2 == 1 else TABLE_ALT_ROW
            _set_cell_fill(data_cell, bg)
            _set_cell_border(data_cell, TABLE_BORDER)

    return table


# ---------------------------------------------------------------------------
# Presentation Factory
# ---------------------------------------------------------------------------

def create_presentation() -> Presentation:
    """Create a blank 16:9 Presentation object with correct dimensions."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


# ---------------------------------------------------------------------------
# Slide Builders
# ---------------------------------------------------------------------------

def _add_blank_slide(prs):
    """Add a blank-layout slide (layout index 6 = Blank)."""
    layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(layout)


def _add_title_bar(slide, title_text: str):
    """Reusable title bar: left-aligned, 28pt bold, TEXT_PRIMARY."""
    add_text_box(
        slide,
        title_text,
        left=MARGIN_LEFT,
        top=TITLE_TOP,
        width=CONTENT_WIDTH,
        height=TITLE_HEIGHT,
        font_size=TITLE_SIZE,
        color=TEXT_PRIMARY,
        bold=True,
        alignment=PP_ALIGN.LEFT,
    )


def _add_title_with_subtitle(slide, title, subtitle, subtitle_color=None):
    """Render title (28pt bold) + subtitle (14pt) below it. Returns bottom position."""
    if subtitle_color is None:
        subtitle_color = TEXT_SECONDARY

    # Title
    add_text_box(
        slide,
        title,
        left=MARGIN_LEFT,
        top=TITLE_TOP,
        width=CONTENT_WIDTH,
        height=Inches(0.55),
        font_size=TITLE_SIZE,
        color=TEXT_PRIMARY,
        bold=True,
        alignment=PP_ALIGN.LEFT,
    )

    # Subtitle just below
    subtitle_top = TITLE_TOP + Inches(0.55)
    add_text_box(
        slide,
        subtitle,
        left=MARGIN_LEFT,
        top=subtitle_top,
        width=CONTENT_WIDTH,
        height=Inches(0.35),
        font_size=Pt(14),
        color=subtitle_color,
        bold=False,
        alignment=PP_ALIGN.LEFT,
    )

    bottom = subtitle_top + Inches(0.4)
    return bottom


def add_callout_box(slide, text, left, top, width, height, accent_color, bg_color=None):
    """
    Draw a callout box with a left accent border and light background.
    Used for 'Key Takeaway' boxes at the bottom of slides.
    """
    if bg_color is None:
        # Very light tint of the accent color
        bg_color = RGBColor(
            min(0xFF, accent_color[0] + (0xFF - accent_color[0]) * 9 // 10),
            min(0xFF, accent_color[1] + (0xFF - accent_color[1]) * 9 // 10),
            min(0xFF, accent_color[2] + (0xFF - accent_color[2]) * 9 // 10),
        )

    # Background rectangle
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb_copy(bg_color)
    bg_shape.line.fill.background()  # no outline

    # Left accent bar (4px = ~50800 EMU wide)
    bar_width = Emu(50800)
    bar_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, bar_width, height
    )
    bar_shape.fill.solid()
    bar_shape.fill.fore_color.rgb = _rgb_copy(accent_color)
    bar_shape.line.fill.background()

    # Text inside the box (indented past accent bar)
    text_left = left + Inches(0.2)
    text_width = width - Inches(0.3)
    add_text_box(
        slide,
        text,
        left=text_left,
        top=top + Inches(0.08),
        width=text_width,
        height=height - Inches(0.16),
        font_size=Pt(13),
        color=TEXT_PRIMARY,
        bold=False,
        alignment=PP_ALIGN.LEFT,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )

    return bg_shape


def add_wave_indicators(slide, waves, top):
    """
    Draw small colored rounded rectangles for each active wave.
    waves: list of tuples [("4G", color), ("5G", color), ("AI", color)]
    Positioned horizontally centered at the given top.
    """
    pill_w = Inches(0.8)
    pill_h = Inches(0.28)
    gap = Inches(0.15)
    total_width = len(waves) * pill_w + (len(waves) - 1) * gap
    start_left = (SLIDE_WIDTH - total_width) // 2

    for i, (label, color) in enumerate(waves):
        x = start_left + i * (pill_w + gap)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, top, pill_w, pill_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb_copy(color)
        shape.line.fill.background()

        add_text_box(
            slide,
            label,
            left=x,
            top=top,
            width=pill_w,
            height=pill_h,
            font_size=Pt(10),
            color=WHITE,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )


def _add_reference_footnote(slide, reference: str):
    """Add a small gray Harvard reference at the bottom-left."""
    add_text_box(
        slide,
        reference,
        left=MARGIN_LEFT,
        top=SLIDE_HEIGHT - Inches(0.5),
        width=Inches(8),
        height=Inches(0.3),
        font_size=SMALL_SIZE,
        color=TEXT_MUTED,
        bold=False,
        alignment=PP_ALIGN.LEFT,
    )


def _add_thin_line(slide, left, top, width, color: RGBColor, height_emu: int = 12700):
    """Draw a thin horizontal line (rectangle shape with minimal height)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Emu(height_emu)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb_copy(color)
    shape.line.fill.background()  # no outline


# ---- 1. Title Slide ----

def create_title_slide(
    prs,
    title: str,
    subtitle: str,
    team_members: list[str],
    module_code: str,
    date: str,
):
    """
    Center-aligned title page.
      - Title 32pt bold
      - Subtitle 18pt
      - Thin blue horizontal line
      - Team members at bottom (12pt)
      - Module + date at very bottom (10pt gray)
    White background, NO gradient.
    """
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)

    # Title — centered
    add_text_box(
        slide,
        title,
        left=Inches(1.5),
        top=Inches(1.5),
        width=SLIDE_WIDTH - Inches(3.0),
        height=Inches(1.0),
        font_size=Pt(32),
        color=TEXT_PRIMARY,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Subtitle — centered
    add_text_box(
        slide,
        subtitle,
        left=Inches(1.5),
        top=Inches(2.5),
        width=SLIDE_WIDTH - Inches(3.0),
        height=Inches(0.6),
        font_size=SUBTITLE_SIZE,
        color=TEXT_SECONDARY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )

    # Thin blue horizontal line
    line_width = Inches(4.0)
    line_left = (SLIDE_WIDTH - line_width) // 2
    _add_thin_line(slide, line_left, Inches(3.25), line_width, ACCENT_BLUE)

    # Team members — one per line, centered
    members_text = "\n".join(team_members)
    add_text_box(
        slide,
        members_text,
        left=Inches(2.0),
        top=Inches(3.6),
        width=SLIDE_WIDTH - Inches(4.0),
        height=Inches(2.5),
        font_size=SMALL_SIZE,
        color=TEXT_PRIMARY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )

    # Module code + date — very bottom
    footer_text = f"{module_code}  |  {date}"
    add_text_box(
        slide,
        footer_text,
        left=Inches(2.0),
        top=SLIDE_HEIGHT - Inches(0.6),
        width=SLIDE_WIDTH - Inches(4.0),
        height=Inches(0.35),
        font_size=SLIDE_NUMBER_SIZE,
        color=TEXT_MUTED,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )

    return slide


# ---- 2. Section Divider ----

def create_section_divider(
    prs,
    task_number: str,
    title: str,
    description: str,
    accent_color: RGBColor = ACCENT_BLUE,
):
    """
    Centered section divider slide.
      - Task number + title (28pt bold)
      - Thin accent bar in wave color
      - Brief description (16pt gray)
    """
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, OFF_WHITE)

    # Task number + title
    heading = f"{task_number}: {title}" if task_number else title
    add_text_box(
        slide,
        heading,
        left=Inches(2.0),
        top=Inches(2.5),
        width=SLIDE_WIDTH - Inches(4.0),
        height=Inches(0.9),
        font_size=TITLE_SIZE,
        color=TEXT_PRIMARY,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Thin accent bar
    bar_width = Inches(2.5)
    bar_left = (SLIDE_WIDTH - bar_width) // 2
    _add_thin_line(slide, bar_left, Inches(3.5), bar_width, accent_color, height_emu=25400)

    # Description
    add_text_box(
        slide,
        description,
        left=Inches(2.0),
        top=Inches(3.9),
        width=SLIDE_WIDTH - Inches(4.0),
        height=Inches(1.0),
        font_size=BODY_SIZE,
        color=TEXT_SECONDARY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )

    return slide


# ---- 3. Content Slide ----

def create_content_slide(
    prs,
    title: str,
    bullets: list[str],
    reference: str | None = None,
    slide_number: int | None = None,
):
    """
    Standard content slide with title + bullet points.
      - Title at top (28pt bold, left-aligned)
      - Bullets using bullet character, 16pt, max 5-6 recommended
      - Optional reference footnote at bottom-left
    Returns the slide object.
    """
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)

    # Title
    _add_title_bar(slide, title)

    # Build bullet text with bullet character
    bullet_lines = "\n".join(f"\u2022  {b}" for b in bullets)

    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT,
        MARGIN_TOP_CONTENT,
        CONTENT_WIDTH,
        CONTENT_HEIGHT - Inches(0.4),  # room for reference
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for idx, bullet in enumerate(bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"\u2022  {bullet}"
        p.font.size = BODY_SIZE
        p.font.color.rgb = _rgb_copy(TEXT_PRIMARY)
        p.font.name = FONT_FAMILY
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)

    if reference:
        _add_reference_footnote(slide, reference)

    if slide_number is not None:
        add_slide_number(slide, slide_number)

    return slide


# ---- 4. Table Slide ----

def create_table_slide(
    prs,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    reference: str | None = None,
    slide_number: int | None = None,
    col_widths: list | None = None,
):
    """
    Slide with title + styled table.
      - Dark header row, alternating data rows
      - Table does NOT fill the entire slide — leaves breathing room
    """
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)

    # Title
    _add_title_bar(slide, title)

    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)

    # Table dimensions — leave margins
    table_left = MARGIN_LEFT
    table_top = MARGIN_TOP_CONTENT + Inches(0.05)
    table_width = CONTENT_WIDTH
    # Height: adaptive but capped so it doesn't fill the slide
    row_height_approx = Inches(0.4)
    table_height = min(
        row_height_approx * num_rows,
        SLIDE_HEIGHT - table_top - Inches(0.9),  # breathing room
    )

    table = create_styled_table(
        slide, num_rows, num_cols, table_left, table_top, table_width, table_height
    )

    # Apply column widths if provided
    if col_widths and len(col_widths) == num_cols:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Populate header
    for col_idx, hdr in enumerate(headers):
        set_cell_text(
            table.cell(0, col_idx),
            hdr,
            font_size=TABLE_HEADER_SIZE,
            color=TABLE_HEADER_TEXT,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

    # Populate data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            set_cell_text(
                table.cell(row_idx + 1, col_idx),
                cell_text,
                font_size=TABLE_BODY_SIZE,
                color=TEXT_PRIMARY,
                bold=False,
                alignment=PP_ALIGN.LEFT,
            )

    if reference:
        _add_reference_footnote(slide, reference)

    if slide_number is not None:
        add_slide_number(slide, slide_number)

    return slide


# ---- 5. Diagram Slide ----

def create_diagram_slide(
    prs,
    title: str,
    image_path: str,
    caption: str | None = None,
    reference: str | None = None,
    slide_number: int | None = None,
):
    """
    Slide with title + centered embedded image + optional caption/reference.
    """
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)

    # Title
    _add_title_bar(slide, title)

    # Image — centered, constrained to content area
    max_img_width = CONTENT_WIDTH - Inches(1.0)
    max_img_height = CONTENT_HEIGHT - Inches(1.0)

    pic = add_image(slide, image_path, Inches(0), Inches(0), width=max_img_width)

    # If the image is taller than allowed, re-add with height constraint
    if pic.height > max_img_height:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = add_image(slide, image_path, Inches(0), Inches(0), height=max_img_height)

    # Center the image horizontally and vertically in content area
    img_left = (SLIDE_WIDTH - pic.width) // 2
    img_top = MARGIN_TOP_CONTENT + (CONTENT_HEIGHT - pic.height - Inches(0.6)) // 2
    pic.left = int(img_left)
    pic.top = int(img_top)

    # Caption below image
    if caption:
        caption_top = pic.top + pic.height + Inches(0.1)
        add_text_box(
            slide,
            caption,
            left=Inches(1.5),
            top=caption_top,
            width=SLIDE_WIDTH - Inches(3.0),
            height=Inches(0.35),
            font_size=SMALL_SIZE,
            color=TEXT_SECONDARY,
            bold=False,
            alignment=PP_ALIGN.CENTER,
        )

    if reference:
        _add_reference_footnote(slide, reference)

    if slide_number is not None:
        add_slide_number(slide, slide_number)

    return slide


# ---- 6. Two-Column Slide ----

def create_two_column_slide(
    prs,
    title: str,
    left_bullets: list[str],
    right_bullets: list[str] | None = None,
    right_image_path: str | None = None,
    reference: str | None = None,
    slide_number: int | None = None,
):
    """
    Title at top, 55/45 split.
      Left column: bullet text.
      Right column: bullets OR image (right_bullets takes precedence).
    """
    slide = _add_blank_slide(prs)
    _set_slide_bg(slide, WHITE)

    # Title
    _add_title_bar(slide, title)

    # Column dimensions (55% / 45% of content width)
    gap = Inches(0.3)
    left_col_width = int(CONTENT_WIDTH * 0.55 - gap // 2)
    right_col_width = int(CONTENT_WIDTH * 0.45 - gap // 2)
    right_col_left = MARGIN_LEFT + left_col_width + gap
    col_height = CONTENT_HEIGHT - Inches(0.4)

    # Left column — bullets
    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP_CONTENT, left_col_width, col_height
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(left_bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"\u2022  {bullet}"
        p.font.size = BODY_SIZE
        p.font.color.rgb = _rgb_copy(TEXT_PRIMARY)
        p.font.name = FONT_FAMILY
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)

    # Right column — bullets or image
    if right_bullets:
        txBox_r = slide.shapes.add_textbox(
            right_col_left, MARGIN_TOP_CONTENT, right_col_width, col_height
        )
        tf_r = txBox_r.text_frame
        tf_r.word_wrap = True
        for idx, bullet in enumerate(right_bullets):
            if idx == 0:
                p = tf_r.paragraphs[0]
            else:
                p = tf_r.add_paragraph()
            p.text = f"\u2022  {bullet}"
            p.font.size = BODY_SIZE
            p.font.color.rgb = _rgb_copy(TEXT_PRIMARY)
            p.font.name = FONT_FAMILY
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(8)

    elif right_image_path and os.path.isfile(right_image_path):
        img_max_w = right_col_width - Inches(0.2)
        img_max_h = col_height - Inches(0.2)
        pic = add_image(
            slide, right_image_path, right_col_left, MARGIN_TOP_CONTENT, width=img_max_w
        )
        # Constrain height
        if pic.height > img_max_h:
            sp = pic._element
            sp.getparent().remove(sp)
            pic = add_image(
                slide, right_image_path, right_col_left, MARGIN_TOP_CONTENT, height=img_max_h
            )
        # Vertically center in column
        pic.top = int(MARGIN_TOP_CONTENT + (col_height - pic.height) // 2)

    if reference:
        _add_reference_footnote(slide, reference)

    if slide_number is not None:
        add_slide_number(slide, slide_number)

    return slide


# ---------------------------------------------------------------------------
# Test Function
# ---------------------------------------------------------------------------

def test_theme():
    """Generate a test presentation with one of each slide type."""
    prs = create_presentation()

    # 1. Title slide
    create_title_slide(
        prs,
        title="Mobilé Inc. — Strategic Analysis 2025",
        subtitle="International Business Strategy with Simulation",
        team_members=[
            "Aayush Man Singh (18029779)",
            "Gaurav Dangol (18029811)",
            "Ruchan Jung Sijapati (20048645)",
            "Shaswat Nibha Maharjan",
            "Bishwesh Ram Shrestha (25030247)",
            "Bishan Subedi (25030193)",
        ],
        module_code="MN7002NI",
        date="10 April 2026",
    )

    # 2. Section divider
    create_section_divider(
        prs,
        task_number="Task 2",
        title="Internal & External Analysis",
        description="VRIO, Value Chain, PESTLE, Porter's Five Forces, Strategic Group, CSF",
        accent_color=WAVE_5G,
    )

    # 3. Content slide
    create_content_slide(
        prs,
        title="5G Market Opportunity",
        bullets=[
            "70%+ urban 5G coverage across all three markets by 2025",
            "Mobilé has NOT launched a 5G handset — critical competitive gap",
            "38% US adoption, 28% EU adoption by end-2024",
            "Asia bifurcated: premium 5G cities + mass 4G rural",
            "First-mover window closing rapidly as rivals ship 5G models",
        ],
        reference="Source: Mobilé Inc. case study, 2025",
        slide_number=7,
    )

    # 4. Table slide
    create_table_slide(
        prs,
        title="VRIO Analysis — Core Resources",
        headers=["Resource", "Value", "Rarity", "Imitability", "Organised", "Outcome"],
        rows=[
            ["Atlanta plant (12 lines)", "Yes", "No", "Easy", "Yes", "Competitive parity"],
            ["Vietnam plant ($6/hr)", "Yes", "Yes", "Hard", "Partial", "Temporary advantage"],
            ["4G revenue ($4.8B)", "Yes", "No", "Easy", "Yes", "Competitive parity"],
            ["Brand (25% share)", "Yes", "No", "Moderate", "Yes", "Competitive parity"],
            ["Cash reserves (>$90M)", "Yes", "No", "N/A", "Constrained", "Threshold resource"],
        ],
        reference="VRIO framework (Barney, 1991)",
        slide_number=9,
        col_widths=[Inches(2.2), Inches(0.8), Inches(0.8), Inches(1.0), Inches(1.1), Inches(2.0)],
    )

    # 5. Diagram slide — skipped unless a test image exists
    test_img = os.path.join(
        os.path.dirname(__file__), "assets", "test_diagram.png"
    )
    if os.path.isfile(test_img):
        create_diagram_slide(
            prs,
            title="Strategic Group Map",
            image_path=test_img,
            caption="Figure 1: Strategic group positioning of Mobilé vs. competitors",
            reference="Adapted from case study data",
            slide_number=17,
        )
    else:
        # Add a placeholder content slide instead
        create_content_slide(
            prs,
            title="Diagram Slide (placeholder — no test image found)",
            bullets=[
                "This slide would embed a centered image (e.g. strategic group map)",
                "Caption appears below the image in 12pt secondary gray",
                f"Looked for: {test_img}",
            ],
            slide_number=17,
        )

    # 6. Two-column slide
    create_two_column_slide(
        prs,
        title="Manufacturing Comparison",
        left_bullets=[
            "Atlanta: 12 lines, $2.9M/line",
            "Labor cost: $28/hr",
            "Proximity to US market",
            "Established 4G production",
        ],
        right_bullets=[
            "Vietnam: expanding, $3.3M/line",
            "Labor cost: $6/hr (78% savings)",
            "Access to Asia supply chain",
            "2-year lead time for new plant",
        ],
        reference="Source: Mobilé Inc. case study, 2025",
        slide_number=12,
    )

    # Save
    output_dir = os.path.join(os.path.dirname(__file__), "output")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "test_theme.pptx")
    prs.save(output_path)
    print(f"Test presentation saved to: {output_path}")
    print(f"  Slides generated: {len(prs.slides)}")


if __name__ == "__main__":
    test_theme()
